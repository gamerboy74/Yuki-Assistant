"""
Tool dispatcher — routes tool calls from the brain to the right handler.

Handles both power tools (powershell, http, file) and plugins.
"""

import json
import re
import urllib.parse
import datetime
from typing import Any
from backend.utils.logger import get_logger

logger = get_logger(__name__)

_UNSAFE_APP_TARGET = re.compile(r"[;&|><`]")


def _looks_unsafe_app_target(target: str) -> bool:
    """Reject obvious shell control characters in app targets."""
    return bool(_UNSAFE_APP_TARGET.search(target))


def _launch_app_target(target: str) -> None:
    """Launch app/URL/protocol on Windows without string-interpolated shell commands."""
    import os
    import subprocess

    lowered = target.lower()

    # Robust URL detection: check for protocol, then check for domain patterns
    is_url = lowered.startswith(("http://", "https://", "www.")) or "://" in target or target.endswith(":")
    
    # If it looks like a domain (e.g. google.com) but lacks protocol
    if not is_url and "." in target and " " not in target and "\\" not in target:
        is_url = True
        if not lowered.startswith("www."):
            target = "https://" + target
        else:
            target = "https://" + target # startfile handles www. but https is safer
    
    if is_url:
        if target.lower().startswith("www."):
            target = "https://" + target
        os.startfile(target)
        return

    # First attempt direct process launch.
    try:
        subprocess.Popen([target], shell=False)
        return
    except Exception:
        pass

    # Fallback for Start menu aliases / registered app names.
    subprocess.Popen(["cmd", "/c", "start", "", target], shell=False)


def _open_chrome_search(query: str) -> str:
    """Open a Google search in Chrome when available, else default browser."""
    import os
    import subprocess

    q = (query or "").strip()
    if not q:
        return "No search query provided."

    encoded = urllib.parse.quote_plus(q)
    url = f"https://www.google.com/search?q={encoded}&tbs=qdr:d"

    try:
        # Prefer direct process launch to avoid cmd metachar parsing of '&' in URLs.
        subprocess.Popen(["chrome", url], shell=False)
        return f"Opened Chrome and searched for '{q}'."
    except Exception:
        # Fallback to registered default browser.
        os.startfile(url)
        return f"Opened browser search for '{q}'."


def _enrich_live_query(query: str) -> str:
    """Bias live sports queries toward fresher authoritative score sources."""
    q = (query or "").strip()
    ql = q.lower()
    if "ipl" in ql and ("score" in ql or "live" in ql or "match" in ql or "life" in ql):
        return f"{q} live score today site:espncricinfo.com OR site:cricbuzz.com"
    return q


def _rank_search_results(query: str, results: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Rank results to prefer current/live and trusted sources for sports/news queries."""
    q = (query or "").lower()
    now_year = str(datetime.datetime.now().year)
    # keep meaningful tokens only
    tokens = [t for t in re.findall(r"[a-z0-9]+", q) if len(t) >= 3]
    trusted = ("espncricinfo", "cricbuzz", "iplt20", "icc-cricket")

    def score(item: dict[str, Any]) -> int:
        title = str(item.get("title", ""))
        body = str(item.get("body", ""))
        href = str(item.get("href", ""))
        blob = f"{title} {body} {href}".lower()

        s = 0
        if "live" in blob:
            s += 8
        if "today" in blob:
            s += 6
        if now_year in blob:
            s += 8
        if any(dom in href.lower() for dom in trusted):
            s += 10

        # reward query token overlap
        for t in tokens:
            if t in blob:
                s += 1

        # penalize clearly old years for current score lookups
        if any(y in blob for y in ("2020", "2021", "2022", "2023", "2024", "2025")) and now_year not in blob:
            s -= 6

        return s

    return sorted(results, key=score, reverse=True)


def dispatch_tool(tool_name: str, arguments: str | dict) -> str:
    """
    Execute a tool call and return the result as a string.
    
    Args:
        tool_name: Function name from the LLM tool call
        arguments: JSON string or dict of parameters
    
    Returns:
        Result string to feed back to the LLM.
    """
    # Parse arguments if string
    if isinstance(arguments, str):
        try:
            params = json.loads(arguments)
        except json.JSONDecodeError:
            return f"Invalid arguments JSON: {arguments[:100]}"
    else:
        params = arguments

    logger.info(f"[DISPATCH] {tool_name}({params})")

    try:
        # ── Power tools ──────────────────────────────────────────────────
        if tool_name == "run_powershell":
            from backend.tools.powershell import run_powershell
            return run_powershell(
                script=params.get("script", ""),
                timeout=params.get("timeout", 10),
            )

        elif tool_name == "http_get":
            from backend.tools.http_request import http_get
            return http_get(
                url=params.get("url", ""),
                timeout=params.get("timeout", 8),
            )

        elif tool_name == "find_file":
            from backend.tools.file_search import find_file
            return find_file(
                name=params.get("name", ""),
                search_dirs=params.get("search_dirs"),
            )

        elif tool_name == "search_internet":
            query = params.get("query", "")
            if not query:
                return "No search query provided."

            search_query = _enrich_live_query(query)

            def _run_search(q):
                # Try ddgs then duckduckgo_search
                try:
                    from ddgs import DDGS
                    return list(DDGS().text(q, max_results=10))
                except Exception:
                    try:
                        from duckduckgo_search import DDGS
                        return list(DDGS().text(q, max_results=10))
                    except Exception as e:
                        raise e

            try:
                results = _run_search(search_query)
                results = _rank_search_results(query, results)[:4]
                if not results:
                    return f"No internet results found for '{query}'."

                output = f"Internet search results for '{query}':\n\n"
                for i, r in enumerate(results, 1):
                    output += f"{i}. {r.get('title')}\n{r.get('body')}\n\n"
                return output
            except Exception as e:
                err_msg = str(e).lower()
                logger.warning(f"search_internet failed: {e}")
                
                # Pivot to Browser if rate-limited (429) or other API blocks
                if "429" in err_msg or "too many requests" in err_msg or "blocked" in err_msg:
                    logger.info(f"Rate limited on API search for '{query}'. Pivoting to Browser Search...")
                    return dispatch_tool("search_in_chrome", {"query": query})
                
                return _search_news_rss(query, limit=5, region="IN")

        elif tool_name == "search_in_chrome":
            query = params.get("query", "")
            if not query: return "No search query provided."

            from backend import chrome_controller
            # Use the new tolerant observant search
            return chrome_controller.search_web_and_read(query)

        elif tool_name == "read_active_tab":
            from backend import chrome_controller
            return chrome_controller.get_page_text(
                max_chars=params.get("max_chars", 2000)
            )

        elif tool_name == "browser_navigate":
            from backend import chrome_controller
            return chrome_controller.navigate(params.get("url", ""))

        elif tool_name == "browser_click":
            from backend import chrome_controller
            return chrome_controller.click_element(params.get("target", ""))

        elif tool_name == "get_page_elements":
            from backend import chrome_controller
            return chrome_controller.get_interactive_elements()

        elif tool_name == "browser_scroll":
            from backend import chrome_controller
            return chrome_controller.scroll(params.get("direction", "down"))

        elif tool_name == "latest_news":
            return _latest_news(
                topic=params.get("topic", "world"),
                limit=params.get("limit", 5),
                region=params.get("region", "IN"),
            )

        elif tool_name == "read_file":
            return _read_file(params.get("path", ""))

        elif tool_name in ("write_file", "design_web_page"):
            return _write_file(
                path=params.get("path", ""),
                content=params.get("content", ""),
                mode=params.get("mode", "overwrite"),
            )

        elif tool_name == "play_youtube":
            # GPT-4o picks the query — we just launch it
            query = params.get("query", "")
            if not query:
                return "No search query provided."
            from backend.plugins.youtube import YouTubePlugin
            result = YouTubePlugin().execute(query=query)
            return result

        elif tool_name == "open_app":
            # Universal App Opener — uses executor.APP_MAP as single source of truth
            from backend.executor import APP_MAP
            name = params.get("name", "").strip().lower()
            if not name:
                return "No app name provided."

            # If the "app" name looks like a file path or has an extension, pivot to open_file
            if "." in name and "\\" in name or ":" in name and "/" not in name:
                return dispatch_tool("open_file", {"path": name})

            target = APP_MAP.get(name, name).strip()

            if not target:
                return "No app name provided."
            if _looks_unsafe_app_target(target):
                return f"Blocked unsafe app target: {name}"

            try:
                _launch_app_target(target)
                return f"Opening {name}..."
            except Exception as e:
                return f"Could not open {name}: {e}"

        elif tool_name == "open_file":
            return _open_file(params.get("path", ""))

        elif tool_name == "type_text":
            text = params.get("text", "")
            if not text:
                return "No text provided to type."
            try:
                import time
                import pyautogui
                # Wait 1.5s to ensure the app window (like Notepad) is focused and ready
                time.sleep(1.5) 
                # Human-like typing interval
                pyautogui.typewrite(text, interval=0.03) 
                return f"Typed successfully."
            except Exception as e:
                return f"Typing failed: {e}"

        # ── Plugin system ────────────────────────────────────────────────
        elif tool_name == "run_plugin":
            from backend.plugins import execute_plugin
            plugin_name = params.pop("plugin_name", "")
            return execute_plugin(plugin_name, params)

        # ── Legacy executor actions (used by fast router) ────────────────
        else:
            # Fall back to the existing executor for fast-router actions
            from backend.executor import execute
            action = {"type": tool_name, "params": params}
            result = execute(action)
            if result is None:
                return "Done."
            elif isinstance(result, dict):
                return result.get("speak", "") or result.get("ui_log", "") or "Done."
            return str(result)

    except Exception as e:
        logger.error(f"[DISPATCH] {tool_name} failed: {e}")
        return f"Tool '{tool_name}' failed: {str(e)[:150]}"


# ── File helpers ──────────────────────────────────────────────────────────────

MAX_FILE_CHARS = 8000  # Cap to avoid token explosion


def _search_news_rss(query: str, limit: int = 5, region: str = "IN") -> str:
    """Fallback internet search using Google News RSS when DDGS is unavailable."""
    import urllib.parse
    import urllib.request
    import xml.etree.ElementTree as ET

    try:
        n = int(limit)
    except Exception:
        n = 5
    n = max(1, min(n, 10))

    q = (query or "").strip()
    if not q:
        return "No search query provided."

    region_clean = (region or "IN").strip().upper()
    query_enc = urllib.parse.quote_plus(q)
    url = (
        "https://news.google.com/rss/search?"
        f"q={query_enc}&hl=en-{region_clean}&gl={region_clean}&ceid={region_clean}:en"
    )

    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=8) as resp:
            data = resp.read()
        root = ET.fromstring(data)
        items = root.findall(".//item")
    except Exception as e:
        return f"Internet search temporarily unavailable for '{q}': {e}"

    if not items:
        return f"No internet results found for '{q}'."

    lines = [f"Internet results (news fallback) for '{q}':"]
    for i, item in enumerate(items[:n], 1):
        title = (item.findtext("title") or "Untitled").strip()
        source = (item.findtext("source") or "Unknown source").strip()
        snippet = (item.findtext("description") or "").strip()
        link = (item.findtext("link") or "").strip()
        lines.append(f"{i}. {title} ({source})")
        if snippet:
            lines.append(snippet)
        if link:
            lines.append(link)
        lines.append("")

    return "\n".join(lines).strip()


def _latest_news(topic: str = "world", limit: int = 5, region: str = "IN") -> str:
    """Fetch latest headlines from Google News RSS (no API key)."""
    import urllib.parse
    import urllib.request
    import xml.etree.ElementTree as ET

    try:
        n = int(limit)
    except Exception:
        n = 5
    n = max(1, min(n, 10))

    topic_clean = (topic or "world").strip()
    region_clean = (region or "IN").strip().upper()

    if topic_clean.lower() in {"world", "global", "latest", "today"}:
        url = f"https://news.google.com/rss?hl=en-{region_clean}&gl={region_clean}&ceid={region_clean}:en"
    else:
        q = urllib.parse.quote_plus(topic_clean)
        url = (
            "https://news.google.com/rss/search?"
            f"q={q}&hl=en-{region_clean}&gl={region_clean}&ceid={region_clean}:en"
        )

    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
    with urllib.request.urlopen(req, timeout=8) as resp:
        data = resp.read()

    root = ET.fromstring(data)
    items = root.findall(".//item")
    if not items:
        return f"No latest headlines found for '{topic_clean}'."

    lines = [f"Latest headlines ({topic_clean}):"]
    for i, item in enumerate(items[:n], 1):
        title = (item.findtext("title") or "Untitled").strip()
        source = (item.findtext("source") or "Unknown source").strip()
        pub = (item.findtext("pubDate") or "").strip()
        link = (item.findtext("link") or "").strip()
        suffix = f" | {pub}" if pub else ""
        lines.append(f"{i}. {title} ({source}){suffix}")
        if link:
            lines.append(f"   {link}")

    return "\n".join(lines)

def _read_file(path: str) -> str:
    """Read and return text content from a file. Supports .txt, .md, .py, .json, .csv, .log, .docx, .pdf."""
    import os
    
    if not path:
        return "No file path provided."
    if not os.path.exists(path):
        return f"File not found: {path}"
    if not os.path.isfile(path):
        return f"Path is a directory, not a file: {path}"
    
    file_size = os.path.getsize(path)
    if file_size > 5 * 1024 * 1024:  # 5 MB cap
        return f"File is too large ({file_size // 1024} KB). Please use a smaller file."
    
    ext = os.path.splitext(path)[1].lower()
    
    try:
        # ── Plain text formats ────────────────────────────────────────────
        if ext in (".txt", ".md", ".py", ".js", ".ts", ".json", ".csv", ".log", ".ini", ".toml", ".yaml", ".yml", ".html", ".xml"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            if len(content) > MAX_FILE_CHARS:
                content = content[:MAX_FILE_CHARS] + f"\n\n[... truncated — file has {len(content)} chars total]"
            return f"Contents of '{os.path.basename(path)}':\n\n{content}"
        
        # ── Word documents (.docx) ────────────────────────────────────────
        elif ext == ".docx":
            try:
                import docx
                doc = docx.Document(path)
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                if len(text) > MAX_FILE_CHARS:
                    text = text[:MAX_FILE_CHARS] + "\n[... truncated]"
                return f"Word document '{os.path.basename(path)}':\n\n{text}"
            except ImportError:
                return "python-docx not installed. Run: pip install python-docx"
        
        # ── PDF documents ─────────────────────────────────────────────────
        elif ext == ".pdf":
            try:
                import pdfplumber
                with pdfplumber.open(path) as pdf:
                    text = "\n".join(
                        page.extract_text() or "" for page in pdf.pages[:15]
                    )
                if len(text) > MAX_FILE_CHARS:
                    text = text[:MAX_FILE_CHARS] + "\n[... truncated, max 8000 chars]"
                return f"PDF '{os.path.basename(path)}':\n\n{text}"
            except ImportError:
                return "pdfplumber not installed. Run: pip install pdfplumber"
        
        # ── PowerPoint presentations (.pptx) ─────────────────────────────
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                from pptx.util import Pt
                
                prs = Presentation(path)
                slides_text: list[str] = []
                
                for i, slide in enumerate(prs.slides, 1):
                    parts: list[str] = [f"--- Slide {i} ---"]
                    
                    # Extract title
                    if slide.shapes.title and slide.shapes.title.text.strip():
                        parts.append(f"Title: {slide.shapes.title.text.strip()}")
                    
                    # Extract all text from shapes (skip title, already captured)
                    body_lines: list[str] = []
                    for shape in slide.shapes:
                        if shape == slide.shapes.title:
                            continue
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                line = para.text.strip()
                                if line:
                                    # Add bullet indent based on paragraph level
                                    indent = "  " * getattr(para, "level", 0)
                                    body_lines.append(f"{indent}• {line}")
                    
                    if body_lines:
                        parts.append("\n".join(body_lines))
                    
                    # Extract speaker notes
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            parts.append(f"[Speaker notes: {notes}]")
                    
                    slides_text.append("\n".join(parts))
                
                full_text = "\n\n".join(slides_text)
                if len(full_text) > MAX_FILE_CHARS:
                    full_text = full_text[:MAX_FILE_CHARS] + f"\n\n[... truncated — {len(prs.slides)} slides total]"
                
                return (
                    f"PowerPoint presentation '{os.path.basename(path)}' "
                    f"({len(prs.slides)} slides):\n\n{full_text}"
                )
            except ImportError:
                return "python-pptx not installed. Run: pip install python-pptx"
        
        else:
            return f"Unsupported file type: {ext}. Supported: .txt .md .py .json .csv .log .docx .pdf .pptx"
    
    except Exception as e:
        return f"Error reading file '{path}': {str(e)[:200]}"


def _open_file(path: str) -> str:
    """Open a file using its default Windows application."""
    import os
    if not path:
        return "No file path provided."
    
    # Check Yuki_Saved first if it's a relative filename
    if not os.path.exists(path):
        from pathlib import Path
        user_home = Path(os.path.expanduser("~")).resolve()
        yuki_saved = user_home / "Documents" / "Yuki_Saved"
        alt_path = yuki_saved / path
        if alt_path.exists():
            path = str(alt_path)
    
    if not os.path.exists(path):
        return f"File not found: {path}"
    
    try:
        os.startfile(path)
        return f"Opened '{os.path.basename(path)}' with its default application."
    except Exception as e:
        return f"Failed to open file: {e}"



def _write_file(path: str, content: str, mode: str = "overwrite") -> str:
    """Write or append content to a file. Auto-launches HTML designs. Hardened for safety."""
    import os
    import webbrowser
    import time
    from pathlib import Path

    # 1. Define standard safe directories
    user_home = Path(os.path.expanduser("~")).resolve()
    yuki_saved = user_home / "Documents" / "Yuki_Saved"
    yuki_designs = user_home / "Documents" / "Yuki_Saved" / "Designs"

    # 2. Path normalization
    if not path:
        # Fallback for empty path: auto-generate in Designs
        yuki_designs.mkdir(parents=True, exist_ok=True)
        path_obj = yuki_designs / f"design_{int(time.time())}.html"
    else:
        # Handle root-relative (e.g., "/landing") or absolute guesses (e.g., "C:\Users\Boss")
        p = Path(path)
        if p.is_absolute() or str(path).startswith(("/", "\\")):
            # If it looks like a system path but isn't in home, or is root-relative,
            # we strip the dangerous parts and move it to Yuki_Saved.
            filename = p.name
            if not filename or filename in (".", ".."):
                filename = f"file_{int(time.time())}.txt"
            path_obj = yuki_saved / filename
        else:
            # Relative path: join with Yuki_Saved
            path_obj = yuki_saved / path

    path_obj = path_obj.resolve()

    if not content:
        return "No content to write."

    try:
        # 3. Final Security Check
        # Ensure we stay within the User's profile (allow Documents, Desktop, etc.)
        if user_home not in path_obj.parents and path_obj != user_home:
            # Emergency fallback: force into Yuki_Saved
            path_obj = yuki_saved / path_obj.name
            path_obj = path_obj.resolve()

        # 4. Create directory and write
        path_obj.parent.mkdir(parents=True, exist_ok=True)

        write_mode = "a" if mode == "append" else "w"
        with open(path_obj, write_mode, encoding="utf-8") as f:
            f.write(content)

        # 5. UI Feedback & Auto-open
        is_html = path_obj.suffix.lower() in (".html", ".htm") or "<html>" in content.lower()[:500]
        
        if is_html:
            # If it's HTML, ensure the extension is correct and open it
            if path_obj.suffix.lower() not in (".html", ".htm"):
                new_path = path_obj.with_suffix(".html")
                os.replace(path_obj, new_path)
                path_obj = new_path
            
            webbrowser.open(f"file:///{path_obj}")
            return f"Design complete! Saved to '{path_obj}' and opened in browser."

        action = "Appended to" if mode == "append" else "Saved"
        return f"{action} '{path_obj.name}' successfully in '{path_obj.parent}'."

    except Exception as e:
        logger.error(f"[DISPATCH] _write_file failed: {e}")
        return f"Error saving file: {str(e)[:150]}"
